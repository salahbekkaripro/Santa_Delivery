"""Generate improved PPTX presentation — Operation Noel V3."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── PALETTE ───────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
C_DARK     = RGBColor(0x16, 0x20, 0x3E)   # deeper navy (section bg)
C_RED      = RGBColor(0xC0, 0x39, 0x2B)   # Santa red
C_GOLD     = RGBColor(0xF3, 0x9C, 0x12)   # Gold
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE     = RGBColor(0x3A, 0x86, 0xFF)   # accent blue
C_GREEN    = RGBColor(0x2E, 0xCC, 0x71)   # success green
C_GREY     = RGBColor(0xEC, 0xF0, 0xF1)   # light grey text
C_ORANGE   = RGBColor(0xE6, 0x7E, 0x22)   # warning orange

W, H = 13.33, 7.5   # slide dimensions in inches

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── PRIMITIVES ────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def text(slide, txt, l, t, w, h, size=16, bold=False, italic=False,
         color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def header(slide, title, icon=""):
    rect(slide, 0, 0, W, H, C_BG)           # background
    rect(slide, 0, 0, W, 1.05, C_RED)        # top bar
    rect(slide, 0, 1.05, W, 0.06, C_GOLD)    # gold line
    label = f"{icon}  {title}" if icon else title
    text(slide, label, 0.35, 0.1, W - 0.5, 0.9,
         size=27, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)


def bullet_slide(prs, title, items, icon=""):
    slide = prs.slides.add_slide(BLANK)
    header(slide, title, icon)
    y = 1.22
    for item in items:
        if isinstance(item, str):
            lvl, txt_, sz, col, bd, gap = 0, item, 16, C_WHITE, False, 0.46
        else:
            lvl = item.get("level", 0)
            txt_ = item.get("text", "")
            sz   = item.get("size", 16 if lvl == 0 else 13)
            col  = item.get("color", C_WHITE if lvl == 0 else C_GREY)
            bd   = item.get("bold", lvl == 0)
            gap  = item.get("gap", 0.46)
        prefix = "▸  " if lvl == 0 else ("    • " if lvl == 1 else "        – ")
        text(slide, prefix + txt_,
             0.4 + lvl * 0.35, y, W - 0.6 - lvl * 0.35, 0.55,
             size=sz, bold=bd, color=col)
        y += gap
    return slide


def section_slide(prs, num, title, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, W, H, C_DARK)
    rect(slide, 0, 0, 0.2, H, C_RED)
    rect(slide, 0.2, 0, W - 0.2, 0.08, C_GOLD)
    text(slide, num,   0.5, 2.3, 2.0, 2.2, size=72, bold=True,
         color=C_RED, align=PP_ALIGN.CENTER)
    text(slide, title, 2.8, 2.7, 10.0, 1.2, size=36, bold=True,
         color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        text(slide, subtitle, 2.8, 3.85, 10.0, 0.7, size=18,
             color=C_GOLD, align=PP_ALIGN.LEFT, italic=True)
    return slide


def two_col(prs, title, left, right, icon=""):
    slide = prs.slides.add_slide(BLANK)
    header(slide, title, icon)
    rect(slide, 6.6, 1.15, 0.05, H - 1.3, C_GOLD)
    y_l = y_r = 1.25
    for items, x, w_ in [(left, 0.4, 6.0), (right, 6.8, 6.2)]:
        y = 1.25
        for b in items:
            if isinstance(b, str):
                t_, sz, bd, col = b, 14, False, C_WHITE
            else:
                t_ = b.get("text", "")
                sz = b.get("size", 14)
                bd = b.get("bold", False)
                col = b.get("color", C_WHITE)
            prefix = "" if t_.startswith("  ") else "▸  "
            text(slide, prefix + t_, x, y, w_, 0.48, size=sz, bold=bd, color=col)
            y += 0.46
    return slide


def metric_slide(prs, title, metrics):
    slide = prs.slides.add_slide(BLANK)
    header(slide, title)
    n = len(metrics)
    bw = (W - 0.9) / n
    for i, m in enumerate(metrics):
        x = 0.45 + i * (bw + 0.05)
        c = m.get("color", C_GOLD)
        rect(slide, x, 1.35, bw - 0.05, 3.8, C_DARK)
        rect(slide, x, 1.35, bw - 0.05, 0.09, c)
        text(slide, m["value"], x, 2.0, bw - 0.05, 1.5,
             size=38, bold=True, color=c, align=PP_ALIGN.CENTER)
        text(slide, m["label"], x, 3.45, bw - 0.05, 0.65,
             size=13, color=C_GREY, align=PP_ALIGN.CENTER)
        if "desc" in m:
            text(slide, m["desc"], x, 4.1, bw - 0.05, 0.8,
                 size=11, color=C_GREY, align=PP_ALIGN.CENTER, italic=True)
    return slide


# ── SLIDE 1 — TITRE ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, C_BG)
rect(slide, 0, 1.7, W, 3.8, C_RED)
rect(slide, 0, 1.65, W, 0.12, C_GOLD)
rect(slide, 0, 5.45, W, 0.12, C_GOLD)
text(slide, "🎄  Operation Noël", 0.5, 1.85, W - 1, 1.4,
     size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
text(slide, "Optimisation des Tournées de Livraison du Père Noël",
     0.5, 3.25, W - 1, 0.75, size=20, color=C_GOLD, align=PP_ALIGN.CENTER)
text(slide, "VRPTW  •  Graphes Réels (OSMnx)  •  OR-Tools  •  Open Data  •  IA Décisionnelle",
     0.5, 4.0, W - 1, 0.5, size=13, color=C_GREY, align=PP_ALIGN.CENTER)
text(slide, "Graphes & Open Data  —  Matière Recherche Opérationnelle  •  2025–2026",
     0.5, 5.75, W - 1, 0.5, size=12, color=C_GREY, align=PP_ALIGN.CENTER)

# ── SLIDE 2 — PLAN ────────────────────────────────────────────────────────────
bullet_slide(prs, "Plan de la Présentation", [
    {"text": "1. Contexte & Problématique — VRPTW, enjeux logistiques", "bold": True, "size": 15, "gap": 0.5},
    {"text": "2. Open Data — OSMnx, Open-Meteo, SRTM ADEME",           "bold": True, "size": 15, "gap": 0.5},
    {"text": "3. Modélisation du Graphe — nœuds, arêtes, poids",        "bold": True, "size": 15, "gap": 0.5},
    {"text": "4. Analyse du Graphe — centralité, robustesse",           "bold": True, "size": 15, "gap": 0.5},
    {"text": "5. Algorithmes — Dijkstra, A*, matrices de coût",         "bold": True, "size": 15, "gap": 0.5},
    {"text": "6. Solveur OR-Tools — phases, profils IA",                "bold": True, "size": 15, "gap": 0.5},
    {"text": "7. Post-traitement — ALNS, ILS, 2-opt, 3-opt",            "bold": True, "size": 15, "gap": 0.5},
    {"text": "8. Incidents & Replanification",                          "bold": True, "size": 15, "gap": 0.5},
    {"text": "9. Résultats & Benchmark — IA vs Humain",                 "bold": True, "size": 15, "gap": 0.5},
    {"text": "10. Architecture & Conclusion",                           "bold": True, "size": 15, "gap": 0.5},
])

# ─────────────────────────────── SECTION 1 ───────────────────────────────────
section_slide(prs, "01", "Contexte & Problématique",
              "Pourquoi optimiser les livraisons du Père Noël ?")

bullet_slide(prs, "Le Défi Logistique", [
    {"text": "Chaque 24 décembre : livraisons massives en temps limité", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.5},
    {"text": "Contraintes réelles du VRPTW :", "size": 16, "bold": True, "gap": 0.38},
    {"text": "Fenêtres de temps : chaque maison a un créneau [earliest, latest]", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Capacité : chaque traîneau porte un poids max de colis", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Météo réelle : neige → vitesse ×0.5, pluie → ×0.77 (Open-Meteo)", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Incidents : rues bloquées → replanification à chaud", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Complexité : NP-difficile", "bold": True, "size": 16, "color": C_RED, "gap": 0.38},
    {"text": "20 clients → 20! ≈ 2.4 × 10¹⁸ combinaisons à explorer", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Résolution exacte impossible → métaheuristiques + OR-Tools", "level": 1, "size": 14, "color": C_RED, "gap": 0.5},
    {"text": "Notre approche : données OSM réelles + Dijkstra + VRPTW + IA apprenante", "size": 16, "bold": True, "color": C_GREEN, "gap": 0.38},
])

# ─────────────────────────────── SECTION 2 ───────────────────────────────────
section_slide(prs, "02", "Open Data",
              "Des données ouvertes qui rendent le problème réaliste")

bullet_slide(prs, "Sources Open Data Utilisées", [
    {"text": "OpenStreetMap via OSMnx", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.4},
    {"text": "Graphe routier réel — Paris 5ème : 124 nœuds, 210 arêtes", "level": 1, "size": 14, "gap": 0.34},
    {"text": "Chaque nœud = intersection ; arête = rue avec longueur + vitesse légale", "level": 1, "size": 14, "gap": 0.34},
    {"text": "Stocké localement : core_data/paris5.graphml (format XML standard)", "level": 1, "size": 14, "gap": 0.44},
    {"text": "Open-Meteo API (gratuite, sans clé)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.4},
    {"text": "Météo en temps réel : température, vent, code météo WMO", "level": 1, "size": 14, "gap": 0.34},
    {"text": "Facteur vitesse dynamique : ☀ ×1.0  •  🌧 ×0.77  •  ❄ ×0.5", "level": 1, "size": 14, "color": C_BLUE, "gap": 0.44},
    {"text": "ADEME Impact CO₂ (facteurs officiels français)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.4},
    {"text": "120 g CO₂/km pour traîneau motorisé (fallback officiel ADEME)", "level": 1, "size": 14, "gap": 0.34},
    {"text": "co2_matrix.npy : matrice précalculée entre tous les points", "level": 1, "size": 14, "gap": 0.44},
    {"text": "OpenTopoData SRTM (relief du terrain)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.4},
    {"text": "Altitude → pente → ajustement du temps et de la consommation", "level": 1, "size": 14, "gap": 0.34},
])

bullet_slide(prs, "Pipeline Open Data : De l'Adresse au Graphe", [
    {"text": "① Requête OSMnx", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": 'ox.graph_from_place("Paris 5e Arrondissement, France", network_type="drive")',
     "level": 1, "size": 13, "color": C_BLUE, "gap": 0.44},
    {"text": "② Enrichissement des arêtes", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "Ajout vitesse légale + facteur météo + pente SRTM → poids final = temps en secondes",
     "level": 1, "size": 14, "gap": 0.44},
    {"text": "③ Calcul des matrices (all-pairs Dijkstra)", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "4 matrices N×N stockées en .npy : temps / distance / CO₂ / risque",
     "level": 1, "size": 14, "gap": 0.44},
    {"text": "④ Matrice composite multicritère", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "Cost = 0.55·t + 0.20·d + 0.15·CO₂ + 0.10·risque",
     "level": 1, "size": 14, "color": C_BLUE, "gap": 0.44},
    {"text": "⑤ Reproductibilité : SHA-256 des matrices → vérification make repro-check",
     "bold": True, "size": 15, "color": C_GREEN, "gap": 0.38},
])

# ─────────────────────────────── SECTION 3 ───────────────────────────────────
section_slide(prs, "03", "Modélisation du Graphe",
              "Un graphe orienté pondéré sur données réelles")

bullet_slide(prs, "Modèle de Graphe Routier", [
    {"text": "Type : Graphe orienté pondéré MultiGraph (OSMnx → NetworkX)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.5},
    {"text": "Orienté : sens unique représenté par arête directionnelle", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Pondéré : poids = temps de trajet (secondes) ou distance (mètres)", "level": 1, "size": 14, "gap": 0.36},
    {"text": "MultiGraph : plusieurs voies parallèles entre deux nœuds possibles", "level": 1, "size": 14, "gap": 0.52},
    {"text": "Caractéristiques — Paris 5ème (core_data/paris5.graphml)", "bold": True, "size": 17, "gap": 0.4},
    {"text": "Nœuds : 124  •  Arêtes : 210", "level": 1, "size": 15, "color": C_GOLD, "gap": 0.36},
    {"text": "Chemin moyen (avg shortest path) : 13.49 sauts", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Diamètre : 35  •  Densité : 0.014  •  Clustering moyen : 0.069", "level": 1, "size": 14, "gap": 0.52},
    {"text": "Vitesses attribuées par type de voie", "bold": True, "size": 16, "gap": 0.38},
    {"text": "residential → 30 km/h  •  tertiary → 50  •  secondary → 60", "level": 1, "size": 14, "color": C_BLUE, "gap": 0.36},
    {"text": "bike → 18 km/h  •  walk → 5 km/h", "level": 1, "size": 14, "color": C_BLUE, "gap": 0.36},
])

# ─────────────────────────────── SECTION 4 ───────────────────────────────────
section_slide(prs, "04", "Analyse du Graphe",
              "Centralité, robustesse et points critiques du réseau")

bullet_slide(prs, "Centralité d'Intermédiarité (Betweenness)", [
    {"text": "BC(v) = fraction des plus courts chemins entre toute paire (s,t) passant par v",
     "bold": True, "size": 16, "color": C_GOLD, "gap": 0.5},
    {"text": "Nœud BC élevé = goulot d'étranglement : sa suppression fragmente le réseau",
     "size": 15, "gap": 0.5},
    {"text": "Top 5 — Paris 5ème :", "bold": True, "size": 16, "gap": 0.38},
    {"text": "#1  Pont Marie                BC = 0.484  →  48% des tournées empruntent ce pont",
     "level": 1, "size": 14, "color": C_RED, "gap": 0.36},
    {"text": "#2  Rue Pont Louis-Philippe   BC = 0.417  →  axe Nord-Sud principal",
     "level": 1, "size": 14, "gap": 0.36},
    {"text": "#3  Rue Saint-Paul            BC = 0.405",
     "level": 1, "size": 14, "gap": 0.36},
    {"text": "#4  Quai de l'Hôtel de Ville  BC = 0.404",
     "level": 1, "size": 14, "gap": 0.36},
    {"text": "#5  Quai de l'Hôtel de Ville  BC = 0.398",
     "level": 1, "size": 14, "gap": 0.52},
    {"text": "Usage : si incident sur Pont Marie → 48% des tournées doivent être reroutées",
     "size": 15, "color": C_GREEN, "gap": 0.38},
])

bullet_slide(prs, "Robustesse du Réseau", [
    {"text": "Test : suppression progressive de nœuds → mesure de la composante connexe",
     "bold": True, "size": 16, "color": C_GOLD, "gap": 0.5},
    {"text": "Attaque ciblée (nœuds à plus fort BC en premier)", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Suppression de 1 nœud  →  perte de 60.5% de la connexité !", "level": 1, "size": 14, "color": C_RED, "gap": 0.36},
    {"text": "Suppression de 5 nœuds →  perte de 66.9%", "level": 1, "size": 14, "gap": 0.36},
    {"text": "Suppression de 10 nœuds → perte de 72.6%", "level": 1, "size": 14, "gap": 0.5},
    {"text": "Attaque aléatoire (nœuds quelconques)", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Suppression de 5 nœuds →  perte de seulement 13.7%", "level": 1, "size": 14, "color": C_GREEN, "gap": 0.5},
    {"text": "Ratio ciblé / aléatoire : ×4.9 — réseau très vulnérable aux attaques ciblées",
     "size": 15, "bold": True, "color": C_GOLD, "gap": 0.5},
    {"text": "Conséquence : incident sur le Pont Marie → replanification immédiate de ~48% des routes",
     "size": 15, "color": C_BLUE, "gap": 0.38},
])

# ─────────────────────────────── SECTION 5 ───────────────────────────────────
section_slide(prs, "05", "Algorithmes de Cheminement",
              "Dijkstra, A* et construction des matrices de coût")

two_col(prs, "Dijkstra vs A* — Comparaison",
    left=[
        {"text": "Dijkstra", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Exploration uniforme (file de priorité)", "size": 14, "color": C_GREY},
        {"text": "  Garantit le chemin optimal", "size": 14, "color": C_GREY},
        {"text": "  Explore TOUS les nœuds accessibles", "size": 14, "color": C_RED},
        {"text": "  Usage : calcul all-pairs → matrices N×N", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "A* (A-étoile)", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Heuristique : distance Haversine (vol d'oiseau)", "size": 14, "color": C_GREY},
        {"text": "  Guide vers la destination → moins de nœuds explorés", "size": 14, "color": C_GREY},
        {"text": "  Optimal si h(n) ≤ distance_réelle (admissible)", "size": 14, "color": C_GREEN},
        {"text": "  Usage : navigation temps réel (un chemin)", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "A* Bidirectionnel", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Deux fronts simultanés : Source→Dest et Dest→Source", "size": 14, "color": C_GREY},
        {"text": "  Convergence quand les fronts se rejoignent", "size": 14, "color": C_GREY},
        {"text": "  -43% de nœuds explorés vs Dijkstra", "size": 14, "color": C_GREEN},
    ],
    right=[
        {"text": "Résultats sur Paris 5ème (124 nœuds)", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Dijkstra :   124 nœuds explorés", "size": 14, "color": C_GREY},
        {"text": "  A*        :   ~80 nœuds  (−35%)", "size": 14, "color": C_GREEN},
        {"text": "  A* bidi   :   ~70 nœuds  (−43%)", "size": 14, "color": C_GREEN},
        {"text": "", "size": 10},
        {"text": "Matrices de coût calculées", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  live_time_matrix.npy   — Temps (s)", "size": 14, "color": C_GREY},
        {"text": "  matrix_5eme.npy        — Distance (m)", "size": 14, "color": C_GREY},
        {"text": "  co2_matrix.npy         — Émissions CO₂", "size": 14, "color": C_GREEN},
        {"text": "  risk_matrix.npy        — Risque incident", "size": 14, "color": C_GREY},
        {"text": "  composite_cost_matrix  — Score multicritère", "size": 14, "color": C_GOLD},
        {"text": "", "size": 10},
        {"text": "case [i, j] = coût Dijkstra du dépôt/client i → client j",
         "size": 13, "color": C_BLUE, "italic": True},
        {"text": "", "size": 10},
        {"text": "Coût composite :", "bold": True, "size": 14, "color": C_GOLD},
        {"text": "  0.55×temps + 0.20×dist + 0.15×CO₂ + 0.10×risque",
         "size": 13, "color": C_BLUE},
    ])

# ─────────────────────────────── SECTION 6 ───────────────────────────────────
section_slide(prs, "06", "Solveur OR-Tools",
              "Résolution VRPTW avec Google OR-Tools + profils IA")

bullet_slide(prs, "Formulation VRPTW", [
    {"text": "Vehicle Routing Problem with Time Windows (VRPTW)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.48},
    {"text": "Entrées :", "bold": True, "size": 16, "gap": 0.36},
    {"text": "Matrice de coût N×N (composite), dépôt unique (nœud 0)", "level": 1, "size": 14, "gap": 0.34},
    {"text": "K traîneaux avec capacité max en kg chacun", "level": 1, "size": 14, "gap": 0.34},
    {"text": "Fenêtres [earliest_i, latest_i] par client + durée de service", "level": 1, "size": 14, "gap": 0.34},
    {"text": "Pénalité drop_penalty par client non livré", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Phases de résolution OR-Tools :", "bold": True, "size": 16, "gap": 0.36},
    {"text": "① Solution initiale : PATH_CHEAPEST_ARC (glouton rapide, quelques ms)", "level": 1, "size": 14, "color": C_BLUE, "gap": 0.34},
    {"text": "② Amélioration : GUIDED_LOCAL_SEARCH (pénalise les arêtes fréquentes)", "level": 1, "size": 14, "color": C_BLUE, "gap": 0.34},
    {"text": "③ Post-traitement : ALNS + ILS + 2-opt + 3-opt + or-opt + double-bridge",
     "level": 1, "size": 14, "color": C_GOLD, "gap": 0.34},
])

bullet_slide(prs, "Profils IA — Politiques du Solveur", [
    {"text": "Les profils IA = politiques de paramétrage d'OR-Tools selon le contexte",
     "bold": True, "size": 16, "color": C_GOLD, "gap": 0.5},
    {"text": "Express", "bold": True, "size": 16, "color": C_BLUE, "gap": 0.36},
    {"text": "PATH_CHEAPEST_ARC + GUIDED_LOCAL_SEARCH, limit 1s → vitesse maximale",
     "level": 1, "size": 14, "gap": 0.44},
    {"text": "Écolo", "bold": True, "size": 16, "color": C_GREEN, "gap": 0.36},
    {"text": "SAVINGS + SIMULATED_ANNEALING → regroupe les trajets, minimise CO₂ et distance",
     "level": 1, "size": 14, "gap": 0.44},
    {"text": "Prudent", "bold": True, "size": 16, "color": C_ORANGE, "gap": 0.36},
    {"text": "Plus grande marge temporelle (time_slack), évite les zones à risque incident",
     "level": 1, "size": 14, "gap": 0.44},
    {"text": "Championne (IA apprenante)", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.36},
    {"text": "Profil appris sur historiques de missions : recommande les meilleurs paramètres",
     "level": 1, "size": 14, "gap": 0.44},
    {"text": "Sélection automatique du nombre de traîneaux K :",
     "bold": True, "size": 15, "color": C_GOLD, "gap": 0.36},
    {"text": "k_min = ceil(poids_total / capacité)  •  score = coût_opérationnel + drop_penalty × non_livrés",
     "level": 1, "size": 13, "color": C_BLUE, "gap": 0.36},
])

# ─────────────────────────────── SECTION 7 ───────────────────────────────────
section_slide(prs, "07", "Post-traitement",
              "ALNS, ILS et amélioration locale de la solution")

bullet_slide(prs, "Amélioration Post-Solve (ALNS + ILS)", [
    {"text": "OR-Tools donne une bonne solution initiale → le post-traitement l'affine",
     "bold": True, "size": 16, "color": C_GOLD, "gap": 0.5},
    {"text": "ALNS — Adaptive Large Neighborhood Search", "bold": True, "size": 16, "gap": 0.36},
    {"text": "Destroy : retirer aléatoirement M clients (Random / Worst-Cost / Relatedness)",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Repair : réinsérer avec heuristique d'insertion au plus bas coût",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Adaptif : les opérateurs efficaces sont récompensés (poids probabiliste)",
     "level": 1, "size": 14, "gap": 0.5},
    {"text": "ILS — Iterated Local Search", "bold": True, "size": 16, "gap": 0.36},
    {"text": "Perturbation double-bridge (échange de segments entre deux tournées)",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Redémarre depuis la meilleure solution connue pour sortir des optima locaux",
     "level": 1, "size": 14, "gap": 0.5},
    {"text": "Micro-optimisations intra-route", "bold": True, "size": 16, "gap": 0.36},
    {"text": "2-opt : échange de deux arêtes dans une même tournée",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "3-opt : échange de trois segments  •  or-opt : déplace un client isolé",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "2-opt* : échange inter-tournées (améliore la répartition entre traîneaux)",
     "level": 1, "size": 14, "gap": 0.34},
])

# ─────────────────────────────── SECTION 8 ───────────────────────────────────
section_slide(prs, "08", "Incidents & Replanification",
              "Reroutage dynamique en cours de mission")

bullet_slide(prs, "Gestion des Incidents en Temps Réel", [
    {"text": "Un incident = rue bloquée, accident ou travaux signalés en cours de mission",
     "bold": True, "size": 16, "color": C_GOLD, "gap": 0.5},
    {"text": "Pipeline de replanification :", "bold": True, "size": 16, "gap": 0.38},
    {"text": "① Signalement via POST /api/missions/{id}/replan-incident",
     "level": 1, "size": 14, "color": C_BLUE, "gap": 0.34},
    {"text": "② Recalcul de la matrice de coût live avec les arêtes bloquées supprimées",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "③ Nouveau passage du solveur OR-Tools sur les clients restants",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "④ Renvoi de la nouvelle tournée au frontend (carte Leaflet mise à jour)",
     "level": 1, "size": 14, "gap": 0.5},
    {"text": "Impact Pont Marie (BC = 0.484) :", "bold": True, "size": 16, "color": C_RED, "gap": 0.38},
    {"text": "~48% des tournées doivent être recalculées si ce nœud est bloqué",
     "level": 1, "size": 14, "color": C_RED, "gap": 0.34},
    {"text": "Cache LRU sur les itinéraires → accélération du recalcul (-60% temps moyen)",
     "level": 1, "size": 14, "color": C_GREEN, "gap": 0.5},
    {"text": "Large scale (≥150 colis) : Génération de tournées candidates + sélection CP-SAT",
     "size": 15, "bold": True, "color": C_GOLD, "gap": 0.36},
])

# ─────────────────────────────── SECTION 9 ───────────────────────────────────
section_slide(prs, "09", "Résultats & Benchmark",
              "IA vs approche humaine sur 5 missions réelles")

metric_slide(prs, "Résultats : IA vs Humain (5 missions, 20 clients)", [
    {"value": "−36.7%", "label": "Réduction du temps",
     "desc": "1 972 s → 1 249 s", "color": C_GOLD},
    {"value": "−36.1%", "label": "Réduction de la distance",
     "desc": "10.86 km → 6.94 km", "color": C_GREEN},
    {"value": "100%",   "label": "Taux de reproductibilité",
     "desc": "SHA-256 identique × 2 passes", "color": C_BLUE},
    {"value": "6",      "label": "Profils testés",
     "desc": "Express / Écolo / Prudent…", "color": C_RED},
])

bullet_slide(prs, "Score Final Multicritère", [
    {"text": "Le score n'est pas seulement le temps : c'est un score pondéré multicritère",
     "bold": True, "size": 17, "color": C_GOLD, "gap": 0.5},
    {"text": "Formule du score composite (poids issus du code) :",
     "bold": True, "size": 16, "gap": 0.38},
    {"text": "Score = 0.55 × temps_normalisé  +  0.20 × distance_normalisée",
     "level": 1, "size": 15, "color": C_BLUE, "gap": 0.34},
    {"text": "           +  0.15 × CO₂_normalisé  +  0.10 × risque_normalisé",
     "level": 1, "size": 15, "color": C_BLUE, "gap": 0.52},
    {"text": "Benchmark : Tournée naïve vs IA optimisée", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Temps : IA gagne 44 minutes sur 5 missions (−36.7%)", "level": 1, "size": 14, "color": C_GREEN, "gap": 0.34},
    {"text": "Distance : 3.92 km de moins par session (−36.1%)", "level": 1, "size": 14, "color": C_GREEN, "gap": 0.34},
    {"text": "CO₂ : proportionnel à la distance → −36% d'émissions estimées", "level": 1, "size": 14, "color": C_GREEN, "gap": 0.52},
    {"text": "Reproductibilité : taux 1.0 — même mission + même seed → résultat identique",
     "size": 15, "bold": True, "color": C_GOLD, "gap": 0.38},
])

# ─────────────────────────────── SECTION 10 ──────────────────────────────────
section_slide(prs, "10", "Architecture & Conclusion",
              "Stack technique et bilan du projet")

two_col(prs, "Architecture Technique",
    left=[
        {"text": "Backend — FastAPI (Python)", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  OR-Tools : solveur VRPTW", "size": 14, "color": C_GREY},
        {"text": "  OSMnx + NetworkX : graphe routier", "size": 14, "color": C_GREY},
        {"text": "  Open-Meteo : météo temps réel", "size": 14, "color": C_GREY},
        {"text": "  SQLite : leaderboard persistant", "size": 14, "color": C_GREY},
        {"text": "  API REST + WebSocket", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "Frontend — Next.js 14 + TypeScript", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Leaflet.js : carte interactive", "size": 14, "color": C_GREY},
        {"text": "  Visualisation des tournées", "size": 14, "color": C_GREY},
        {"text": "  Mode Versus Humain vs IA", "size": 14, "color": C_GREY},
        {"text": "  Replay animé + debriefing", "size": 14, "color": C_GREY},
    ],
    right=[
        {"text": "Fonctionnalités clés", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Mode Joueur : l'humain choisit ses routes", "size": 14, "color": C_GREY},
        {"text": "  Mode IA : OR-Tools résout seul", "size": 14, "color": C_GREY},
        {"text": "  Mode Versus : comparaison temps réel", "size": 14, "color": C_GREY},
        {"text": "  Auto-suggestion : IA conseille le joueur", "size": 14, "color": C_GREY},
        {"text": "  Incidents : reroutage automatique", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "DevOps & Qualité", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Docker Compose (make docker)", "size": 14, "color": C_GREY},
        {"text": "  pytest : 15+ tests automatisés", "size": 14, "color": C_GREY},
        {"text": "  make repro-check : SHA-256", "size": 14, "color": C_GREEN},
        {"text": "  Déterminisme garanti (seed explicite)", "size": 14, "color": C_GREEN},
    ])

bullet_slide(prs, "Ce Que Nous Avons Accompli", [
    {"text": "Graphes & Open Data", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Graphe réel OSMnx (124 nœuds, 210 arêtes) + analyse centralité + robustesse",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Météo live (Open-Meteo) + relief SRTM + CO₂ ADEME + adresses géocodées",
     "level": 1, "size": 14, "gap": 0.48},
    {"text": "Algorithmique", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Dijkstra + A* + A* bidirectionnel → matrices de coût N×N multicritères",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "VRPTW via OR-Tools + GLS + SA + Tabu + ALNS + ILS + 2-opt + 3-opt",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "K-Means spatial custom (NumPy) pour la sectorisation large scale",
     "level": 1, "size": 14, "gap": 0.48},
    {"text": "Résultats mesurés", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "−36.7% de temps  •  −36.1% de distance vs approche humaine",
     "level": 1, "size": 14, "color": C_GREEN, "gap": 0.34},
    {"text": "Reproductibilité 100%  •  déterminisme garanti sur toutes les politiques",
     "level": 1, "size": 14, "color": C_GREEN, "gap": 0.34},
])

bullet_slide(prs, "Limites & Perspectives", [
    {"text": "Limites actuelles", "bold": True, "size": 17, "color": C_RED, "gap": 0.44},
    {"text": "Simulation : les fenêtres de temps et incidents sont générés, pas saisis en vrai",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Optimalité globale non garantie (NP-difficile → approximations)",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "OSMnx lent sur grands périmètres (>500m de rayon)",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "IA apprenante pas encore assez entraînée pour battre les profils fixes",
     "level": 1, "size": 14, "gap": 0.5},
    {"text": "Perspectives", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.44},
    {"text": "Données GTFS Île-de-France → graphe multimodal (véhicule + métro)",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Trafic temps réel (HERE / TomTom) pour pondérer les arêtes dynamiquement",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "Renforcement Learning pour remplacer l'heuristique gloutonne initiale",
     "level": 1, "size": 14, "gap": 0.34},
    {"text": "PageRank + Eigenvector Centrality sur le graphe routier",
     "level": 1, "size": 14, "gap": 0.34},
])

# ── SLIDE FINALE ──────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, C_BG)
rect(slide, 0, 2.4, W, 3.0, C_RED)
rect(slide, 0, 2.35, W, 0.12, C_GOLD)
rect(slide, 0, 5.35, W, 0.12, C_GOLD)
text(slide, "Merci  🎄", 0.5, 2.6, W - 1, 1.2,
     size=58, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
text(slide, "Des questions ?",
     0.5, 3.85, W - 1, 0.7, size=22, color=C_GOLD, align=PP_ALIGN.CENTER)
text(slide, "Operation Noël  •  VRPTW + Graphes Réels + OR-Tools + Open Data  •  2025–2026",
     0.5, 6.05, W - 1, 0.5, size=12, color=C_GREY, align=PP_ALIGN.CENTER)

# ── SAVE ──────────────────────────────────────────────────────────────────────
output = "/home/bekkari/Documents/Graphes/Noel/Operation_Noel_V3.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
