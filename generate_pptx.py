"""Generate PPTX presentation for Operation Noel project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color palette
C_BG_DARK = RGBColor(0x0D, 0x1B, 0x2A)       # Dark navy
C_RED = RGBColor(0xC0, 0x39, 0x2B)            # Santa red
C_GOLD = RGBColor(0xF3, 0x9C, 0x12)           # Gold
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE = RGBColor(0x3A, 0x86, 0xFF)
C_GREEN = RGBColor(0x2E, 0xCC, 0x71)
C_GREY = RGBColor(0xEC, 0xF0, 0xF1)
C_DARK_SECTION = RGBColor(0x16, 0x20, 0x3E)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def add_rect(slide, left, top, width, height, fill_color, transparency=0):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, font_size=18,
             bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_slide(prs, title, bullets, icon=""):
    """Standard content slide with title + bullet list."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    # Background
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG_DARK)
    # Red top bar
    add_rect(slide, 0, 0, 13.33, 1.1, C_RED)
    # Title
    add_text(slide, f"{icon}  {title}" if icon else title,
             0.3, 0.12, 12.7, 0.9, font_size=28, bold=True,
             color=C_WHITE, align=PP_ALIGN.LEFT)
    # Gold accent line
    add_rect(slide, 0.3, 1.15, 12.5, 0.05, C_GOLD)

    # Bullets
    y = 1.4
    for bullet in bullets:
        if isinstance(bullet, dict):
            level = bullet.get("level", 0)
            text = bullet.get("text", "")
            size = bullet.get("size", 17 if level == 0 else 14)
            col = bullet.get("color", C_WHITE if level == 0 else C_GREY)
            prefix = "▸  " if level == 0 else "    • "
            bold = bullet.get("bold", level == 0)
        else:
            level, text, size, col, bold = 0, bullet, 17, C_WHITE, False
            prefix = "▸  "
        add_text(slide, prefix + text, 0.4 + level * 0.3,
                 y, 12.5 - level * 0.3, 0.5, font_size=size,
                 bold=bold, color=col)
        y += bullet.get("gap", 0.48) if isinstance(bullet, dict) else 0.48
    return slide


def add_title_slide(prs):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG_DARK)
    # Big red block
    add_rect(slide, 0, 1.8, 13.33, 3.6, C_RED)
    # Decorative gold stripe
    add_rect(slide, 0, 1.75, 13.33, 0.12, C_GOLD)
    add_rect(slide, 0, 5.4, 13.33, 0.12, C_GOLD)

    add_text(slide, "Operation Noël", 0.5, 2.0, 12.3, 1.4,
             font_size=56, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Optimisation Logistique Urbaine du Père Noël",
             0.5, 3.4, 12.3, 0.7, font_size=22, bold=False,
             color=C_GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "Vehicle Routing Problem with Time Windows (VRPTW) • Open Data • Théorie des Graphes",
             0.5, 4.05, 12.3, 0.5, font_size=14, bold=False,
             color=C_GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "Matière : Graphes & Open Data  |  2025–2026",
             0.5, 5.7, 12.3, 0.5, font_size=13, bold=False,
             color=C_GREY, align=PP_ALIGN.CENTER)
    return slide


def add_section_divider(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, C_DARK_SECTION)
    add_rect(slide, 0, 0, 0.18, 7.5, C_RED)
    add_rect(slide, 0.18, 0, 13.15, 0.08, C_GOLD)

    add_text(slide, number, 0.5, 2.5, 2.0, 2.0, font_size=72, bold=True,
             color=C_RED, align=PP_ALIGN.CENTER)
    add_text(slide, title, 2.6, 2.8, 10.0, 1.2, font_size=36, bold=True,
             color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 2.6, 3.9, 10.0, 0.7, font_size=18,
                 color=C_GOLD, align=PP_ALIGN.LEFT, italic=True)
    return slide


def add_two_col_slide(prs, title, left_bullets, right_bullets, icon=""):
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 1.1, C_RED)
    add_text(slide, f"{icon}  {title}" if icon else title,
             0.3, 0.12, 12.7, 0.9, font_size=28, bold=True, color=C_WHITE)
    add_rect(slide, 0.3, 1.15, 12.5, 0.05, C_GOLD)

    # Separator
    add_rect(slide, 6.6, 1.3, 0.04, 5.8, C_GOLD)

    # Left column
    y = 1.4
    for b in left_bullets:
        text = b if isinstance(b, str) else b.get("text", "")
        size = 15 if isinstance(b, str) else b.get("size", 15)
        bold = False if isinstance(b, str) else b.get("bold", False)
        col = C_WHITE if isinstance(b, str) else b.get("color", C_WHITE)
        prefix = "▸  " if not text.startswith("  ") else ""
        add_text(slide, prefix + text, 0.4, y, 6.0, 0.45, font_size=size,
                 bold=bold, color=col)
        y += 0.46

    # Right column
    y = 1.4
    for b in right_bullets:
        text = b if isinstance(b, str) else b.get("text", "")
        size = 15 if isinstance(b, str) else b.get("size", 15)
        bold = False if isinstance(b, str) else b.get("bold", False)
        col = C_WHITE if isinstance(b, str) else b.get("color", C_WHITE)
        prefix = "▸  " if not text.startswith("  ") else ""
        add_text(slide, prefix + text, 6.8, y, 6.2, 0.45, font_size=size,
                 bold=bold, color=col)
        y += 0.46
    return slide


def add_metric_slide(prs, title, metrics):
    """Slide with big metric boxes."""
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(slide, 0, 0, 13.33, 7.5, C_BG_DARK)
    add_rect(slide, 0, 0, 13.33, 1.1, C_RED)
    add_text(slide, title, 0.3, 0.12, 12.7, 0.9, font_size=28,
             bold=True, color=C_WHITE)
    add_rect(slide, 0.3, 1.15, 12.5, 0.05, C_GOLD)

    box_w = 12.4 / len(metrics)
    for i, m in enumerate(metrics):
        x = 0.45 + i * (box_w + 0.1)
        add_rect(slide, x, 1.5, box_w, 3.5, C_DARK_SECTION)
        add_rect(slide, x, 1.5, box_w, 0.08, m.get("color", C_GOLD))
        add_text(slide, m["value"], x, 2.1, box_w, 1.5, font_size=40,
                 bold=True, color=m.get("color", C_GOLD), align=PP_ALIGN.CENTER)
        add_text(slide, m["label"], x, 3.5, box_w, 0.6, font_size=13,
                 bold=False, color=C_GREY, align=PP_ALIGN.CENTER)
        if "desc" in m:
            add_text(slide, m["desc"], x, 4.1, box_w, 0.7, font_size=11,
                     color=C_GREY, align=PP_ALIGN.CENTER, italic=True)

    return slide


# ── BUILD SLIDES ──────────────────────────────────────────────────────────────

# 1. Title
add_title_slide(prs)

# 2. Plan de la présentation
add_bullet_slide(prs, "Plan de la Présentation", [
    {"text": "Contexte & Problématique", "bold": True, "size": 18},
    {"text": "Open Data : sources et données utilisées", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Théorie des Graphes : modélisation du réseau routier", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Analyse du graphe (centralité, robustesse)", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Algorithmes de cheminement (Dijkstra, A*, A* bidirectionnel)", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Problème VRPTW & optimisation (OR-Tools)", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Clustering spatial (K-Means)", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Profils IA & résultats", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Architecture & démo", "bold": True, "size": 18, "gap": 0.52},
    {"text": "Conclusion & perspectives", "bold": True, "size": 18, "gap": 0.52},
], icon="")

# ── SECTION 1 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "01", "Contexte & Problématique",
                    "Pourquoi optimiser les livraisons du Père Noël ?")

add_bullet_slide(prs, "Le défi logistique", [
    {"text": "Chaque 24 décembre, des milliards de colis doivent être livrés en une nuit", "size": 17, "gap": 0.55},
    {"text": "Contraintes réelles imposées :", "size": 17, "bold": True, "gap": 0.42},
    {"text": "Fenêtres de temps : chaque maison a un créneau de livraison", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Capacité limitée : chaque traîneau ne peut porter qu'un certain poids", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Conditions météo : neige, vent → vitesse réduite (API Open-Meteo)", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Incidents routiers : rues bloquées, travaux, accidents", "level": 1, "size": 15, "gap": 0.42},
    {"text": "Ce problème = VRPTW (Vehicle Routing Problem with Time Windows)", "size": 17, "bold": True, "color": C_GOLD, "gap": 0.55},
    {"text": "Complexité NP-difficile → impossible à résoudre par force brute au-delà de ~20 clients", "size": 15, "gap": 0.42},
    {"text": "Notre solution : graphes réels (OSMnx) + solveur (OR-Tools) + IA apprenante", "size": 16, "bold": True, "color": C_GREEN, "gap": 0.42},
], icon="")

# ── SECTION 2 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "02", "Open Data",
                    "Les sources de données ouvertes au cœur du projet")

add_bullet_slide(prs, "Sources Open Data utilisées", [
    {"text": "OpenStreetMap (OSMnx)", "bold": True, "size": 18, "color": C_GOLD, "gap": 0.4},
    {"text": "Graphe routier réel de Paris 5ème : 124 nœuds, 210 arêtes", "level": 1, "size": 15, "gap": 0.36},
    {"text": "Chaque nœud = intersection, chaque arête = segment de rue avec longueur + vitesse légale", "level": 1, "size": 15, "gap": 0.36},
    {"text": "Exporté en .graphml via OSMnx → persisté localement (core_data/paris5.graphml)", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Open-Meteo API", "bold": True, "size": 18, "color": C_GOLD, "gap": 0.4},
    {"text": "Météo en temps réel (température, vitesse du vent, code météo)", "level": 1, "size": 15, "gap": 0.36},
    {"text": "Impact dynamique sur la vitesse des traîneaux (ex. neige → -30%)", "level": 1, "size": 15, "gap": 0.36},
    {"text": "Gratuite, sans clé API, haute résolution géographique", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Données de livraison générées (synthetic open data)", "bold": True, "size": 18, "color": C_GOLD, "gap": 0.4},
    {"text": "Adresses réelles géocodées via OSMnx (pas de simulation fictive)", "level": 1, "size": 15, "gap": 0.36},
    {"text": "Matrices de coût temps stockées en .npy (NumPy binaire)", "level": 1, "size": 15, "gap": 0.36},
], icon="")

add_bullet_slide(prs, "Pipeline Open Data — De la rue au graphe", [
    {"text": "1. Requête OSMnx", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": 'osmnx.graph_from_place("Paris 5e Arrondissement, France", network_type="drive")', "level": 1, "size": 13, "color": C_LIGHT_BLUE, "gap": 0.42},
    {"text": "2. Projection & nettoyage", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "Simplification du graphe, suppression des nœuds isolés, projection UTM", "level": 1, "size": 14, "gap": 0.42},
    {"text": "3. Calcul des poids temporels", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "Pour chaque arête : t = distance / (speed_limit × weather_factor)", "level": 1, "size": 13, "color": C_LIGHT_BLUE, "gap": 0.42},
    {"text": "4. Matrice de coût (all-pairs Dijkstra)", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "Précalcul de toutes les distances entre points de livraison → matrice N×N stockée en .npy", "level": 1, "size": 14, "gap": 0.42},
    {"text": "5. Reproductibilité garantie", "bold": True, "size": 16, "color": C_GOLD, "gap": 0.38},
    {"text": "Hash SHA-256 du graphe + matrices → vérification make repro-check", "level": 1, "size": 14, "gap": 0.38},
], icon="")

# ── SECTION 3 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "03", "Théorie des Graphes",
                    "Modélisation mathématique du réseau routier parisien")

add_bullet_slide(prs, "Modélisation du Réseau Routier", [
    {"text": "Type : Graphe orienté pondéré (Directed Weighted MultiGraph)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.52},
    {"text": "Orienté : les rues à sens unique sont représentées avec des arêtes directionnelles", "level": 1, "size": 15, "gap": 0.42},
    {"text": "Pondéré : chaque arête a un poids = temps de trajet en secondes", "level": 1, "size": 15, "gap": 0.42},
    {"text": "Multi-graphe : deux nœuds peuvent être reliés par plusieurs arêtes (voies parallèles)", "level": 1, "size": 15, "gap": 0.55},
    {"text": "Caractéristiques du graphe (Paris 5ème)", "bold": True, "size": 17, "gap": 0.42},
    {"text": "Nœuds : 124  |  Arêtes : 210", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Chemin moyen (avg shortest path) : 13.49 sauts", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Diamètre du graphe : 35", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Densité : 0.014  |  Coefficient de clustering moyen : 0.069", "level": 1, "size": 15, "gap": 0.55},
    {"text": "Format de stockage : .graphml (standard XML interopérable)", "size": 15, "color": C_LIGHT_BLUE, "gap": 0.38},
], icon="")

# ── SECTION 4 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "04", "Analyse du Graphe",
                    "Centralité, robustesse et points critiques du réseau")

add_bullet_slide(prs, "Analyse de Centralité d'Intermédiarité", [
    {"text": "Betweenness Centrality (BC) : fraction des plus courts chemins passant par un nœud", "size": 16, "bold": True, "color": C_GOLD, "gap": 0.52},
    {"text": "Un BC élevé = nœud critique → sa suppression fragmente le réseau", "size": 15, "gap": 0.52},
    {"text": "Top 5 nœuds critiques — Paris 5ème", "bold": True, "size": 17, "gap": 0.42},
    {"text": "#1  Pont Marie               (BC = 0.484)  → ~48% des chemins", "level": 1, "size": 14, "color": C_RED, "gap": 0.38},
    {"text": "#2  Rue du Pont Louis-Philippe (BC = 0.417)  → axe Nord-Sud majeur", "level": 1, "size": 14, "gap": 0.38},
    {"text": "#3  Rue Saint-Paul             (BC = 0.405)  → entrée de quartier", "level": 1, "size": 14, "gap": 0.38},
    {"text": "#4  Quai de l'Hôtel de Ville  (BC = 0.404)  → pont vers Île de la Cité", "level": 1, "size": 14, "gap": 0.38},
    {"text": "#5  Quai de l'Hôtel de Ville  (BC = 0.398)  → nœud adjacent", "level": 1, "size": 14, "gap": 0.52},
    {"text": "Application : l'IA reroute automatiquement si un nœud critique est en incident", "size": 15, "color": C_GREEN, "gap": 0.38},
], icon="")

add_bullet_slide(prs, "Analyse de Robustesse du Réseau", [
    {"text": "Simulation : suppression progressive de nœuds → mesure de la perte de connectivité", "size": 16, "bold": True, "color": C_GOLD, "gap": 0.52},
    {"text": "Attaque ciblée (nœuds les plus centraux en premier)", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Suppression de 1 nœud  →  perte de 60.5% de la composante connexe !", "level": 1, "size": 15, "color": C_RED, "gap": 0.38},
    {"text": "Suppression de 5 nœuds →  perte de 66.9%", "level": 1, "size": 15, "gap": 0.38},
    {"text": "Suppression de 10 nœuds → perte de 72.6%", "level": 1, "size": 15, "gap": 0.52},
    {"text": "Attaque aléatoire (nœuds quelconques)", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Suppression de 5 nœuds →  perte de seulement 13.7%", "level": 1, "size": 15, "color": C_GREEN, "gap": 0.52},
    {"text": "Ratio de dévastation : 4.9×  →  le réseau est très vulnérable aux attaques ciblées", "size": 15, "bold": True, "color": C_GOLD, "gap": 0.52},
    {"text": "Implication : un seul incident sur le Pont Marie force >48% des tournées à se recalculer", "size": 15, "color": C_LIGHT_BLUE, "gap": 0.38},
], icon="")

# ── SECTION 5 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "05", "Algorithmes de Cheminement",
                    "Dijkstra, A* et A* bidirectionnel sur graphe routier réel")

add_two_col_slide(prs, "Comparaison des Algorithmes de Cheminement",
    left_bullets=[
        {"text": "Dijkstra", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Exploration uniforme depuis la source", "size": 14, "color": C_GREY},
        {"text": "  Garantit le chemin optimal", "size": 14, "color": C_GREY},
        {"text": "  Explore TOUS les nœuds accessibles", "size": 14, "color": C_RED},
        {"text": "  Usage : calcul de la matrice de coût globale", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "A* (A-Star)", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Heuristique : distance Haversine (vol d'oiseau)", "size": 14, "color": C_GREY},
        {"text": "  Guide l'exploration vers la destination", "size": 14, "color": C_GREY},
        {"text": "  Optimal si heuristique admissible", "size": 14, "color": C_GREEN},
        {"text": "  Usage : navigation en temps réel", "size": 14, "color": C_GREY},
    ],
    right_bullets=[
        {"text": "A* Bidirectionnel", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Deux recherches simultanées :", "size": 14, "color": C_GREY},
        {"text": "  Source → Destination ET Dest → Source", "size": 14, "color": C_LIGHT_BLUE},
        {"text": "  Convergence quand les deux fronts se rejoignent", "size": 14, "color": C_GREY},
        {"text": "  Explore 20-40% de nœuds en moins que Dijkstra", "size": 14, "color": C_GREEN},
        {"text": "  Même résultat optimal, plus rapide", "size": 14, "color": C_GREEN},
        {"text": "", "size": 10},
        {"text": "Résultat sur le graphe Paris 5ème", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Dijkstra : 124 nœuds explorés en moyenne", "size": 14, "color": C_GREY},
        {"text": "  A*        : ~80 nœuds (35% de gain)", "size": 14, "color": C_GREEN},
        {"text": "  A* bidi   : ~70 nœuds (43% de gain)", "size": 14, "color": C_GREEN},
    ])

# ── SECTION 6 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "06", "VRPTW & OR-Tools",
                    "Optimisation combinatoire avec Google OR-Tools")

add_bullet_slide(prs, "Le Problème VRPTW", [
    {"text": "Vehicle Routing Problem with Time Windows", "bold": True, "size": 18, "color": C_GOLD, "gap": 0.5},
    {"text": "Objectif : minimiser le temps total de livraison sous contraintes", "size": 16, "gap": 0.45},
    {"text": "Contraintes modélisées", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Fenêtres de temps [earliest, latest] par client", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Capacité maximale par traîneau (poids total des colis)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Retour au dépôt obligatoire après chaque tournée", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Pénalité pour abandon de client (drop_penalty)", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Complexité : NP-difficile", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Pour N=20 clients → espace de solutions : 20! ≈ 2.4 × 10¹⁸ combinaisons", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Résolution exacte impossible → métaheuristiques nécessaires", "level": 1, "size": 15, "color": C_RED, "gap": 0.38},
], icon="")

add_bullet_slide(prs, "Stratégie de Résolution OR-Tools", [
    {"text": "Phase 1 : Solution initiale (PATH_CHEAPEST_ARC)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Heuristique gloutonne : connecte les arêtes les moins coûteuses en priorité", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Rapide, solution de départ réalisable en quelques ms", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Phase 2 : Amélioration (GUIDED_LOCAL_SEARCH)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Explore les voisinages de la solution courante", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Pénalise les arêtes fréquemment utilisées → évite les optima locaux", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Limite de temps configurable (1s à 30s selon contexte)", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Phase 3 : Post-traitement (ALNS + ILS)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "ALNS : Adaptive Large Neighborhood Search", "level": 1, "size": 15, "gap": 0.35},
    {"text": "ILS : Iterated Local Search avec garde d'intégrité", "level": 1, "size": 15, "gap": 0.35},
], icon="")

add_bullet_slide(prs, "Politiques OR-Tools testées", [
    {"text": "6 politiques évaluées sur les mêmes instances", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.48},
    {"text": "pca_gls_fast     → PATH_CHEAPEST_ARC + GLS, rapide (1s)", "level": 1, "size": 15, "color": C_LIGHT_BLUE, "gap": 0.38},
    {"text": "pca_sa_balanced  → PATH_CHEAPEST_ARC + Simulated Annealing, équilibré", "level": 1, "size": 15, "gap": 0.38},
    {"text": "pci_gls_deep     → PATH_CHEAPEST_INSERT + GLS, approfondi (30s)", "level": 1, "size": 15, "gap": 0.38},
    {"text": "savings_tabu     → Savings heuristic + Tabu Search", "level": 1, "size": 15, "gap": 0.38},
    {"text": "pca_gls_distance → Optimisation distance (profil Ecolo)", "level": 1, "size": 15, "color": C_GREEN, "gap": 0.38},
    {"text": "pci_gls_distance → Insertion + distance (profil Ecolo approfondi)", "level": 1, "size": 15, "color": C_GREEN, "gap": 0.52},
    {"text": "Reproductibilité : toutes les politiques → taux 1.0 (100%)", "bold": True, "size": 16, "color": C_GREEN, "gap": 0.48},
    {"text": "Déterminisme garanti par seed explicite propagée à OR-Tools, ALNS et ILS", "size": 14, "level": 1, "gap": 0.38},
], icon="")

# ── SECTION 7 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "07", "Clustering Spatial",
                    "Sectorisation K-Means pour la scalabilité")

add_bullet_slide(prs, "Sectorisation par K-Means", [
    {"text": "Problème : le VRPTW est NP-difficile → temps exponentiel avec N clients", "size": 16, "bold": True, "color": C_GOLD, "gap": 0.48},
    {"text": "Solution : Divide & Conquer en 2 étapes", "size": 16, "bold": True, "gap": 0.38},
    {"text": "Étape 1 — Clustering : diviser les clients en K secteurs (1 par traîneau)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Étape 2 — Optimisation locale : résoudre un TSP dans chaque secteur", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Algorithme K-Means (implémentation NumPy, sans scikit-learn)", "size": 16, "bold": True, "color": C_GOLD, "gap": 0.38},
    {"text": "Initialisation : K centroïdes aléatoires (ou K-Means++)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Attribution : chaque client → secteur dont le centroïde est le plus proche", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Mise à jour : centroïde = moyenne des positions du secteur", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Convergence : répéter jusqu'à stabilité des affectations", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Résultat Paris 5ème : Zone 0 (N-O: 2 clients), Zone 1 (Est: 5), Zone 2 (S-E: 3)", "size": 15, "color": C_GREEN, "gap": 0.38},
    {"text": "Carte interactive générée : production_output/clustering_map.html", "size": 14, "color": C_LIGHT_BLUE, "gap": 0.38},
], icon="")

# ── SECTION 8 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "08", "Profils IA & Résultats",
                    "L'IA qui apprend à choisir la meilleure stratégie")

add_bullet_slide(prs, "Profils IA et Auto-Suggestion", [
    {"text": "6 profils IA disponibles", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Express      → Minimise le temps de trajet (vitesse max)", "level": 1, "size": 15, "color": C_LIGHT_BLUE, "gap": 0.35},
    {"text": "Ecolo        → Minimise la distance et le CO₂", "level": 1, "size": 15, "color": C_GREEN, "gap": 0.35},
    {"text": "Prudent      → Évite les zones à risque incident", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Opportuniste → Optimise le score pondéré multi-objectif", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Agressive    → Maximum de colis livrés, même hors délai", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Championne   → Profil appris par le modèle IA sur données historiques", "level": 1, "size": 15, "color": C_GOLD, "gap": 0.52},
    {"text": "IA apprenante (modèle v2.0)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.42},
    {"text": "Recommande le meilleur profil selon le contexte (taille, météo, incidents)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Auto-Tuner OR-Tools : apprend les hyperparamètres optimaux du solveur", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Split stratifié par contexte pour une évaluation fiable", "level": 1, "size": 15, "gap": 0.35},
], icon="")

add_metric_slide(prs, "Résultats : IA vs Humain (5 missions, 20 clients)", [
    {"value": "36.7%", "label": "Réduction du temps de trajet",
     "desc": "1972s (humain) → 1249s (IA)", "color": C_GOLD},
    {"value": "36.1%", "label": "Réduction de la distance",
     "desc": "10.86km → 6.94km", "color": C_GREEN},
    {"value": "53.2%", "label": "Score benchmark (économies)",
     "desc": "44min de trajet économisées", "color": C_LIGHT_BLUE},
    {"value": "1.0", "label": "Taux de reproductibilité",
     "desc": "SHA-256 identique sur 2 passes", "color": C_RED},
])

add_bullet_slide(prs, "Benchmark Multi-Villes", [
    {"text": "Script multi_city_benchmark.py : tester plusieurs villes en un seul run", "size": 16, "bold": True, "color": C_GOLD, "gap": 0.52},
    {"text": "Villes supportées (données OSMnx)", "bold": True, "size": 16, "gap": 0.38},
    {"text": "Le Marais, Paris  |  Mitte, Berlin  |  Vieux Lyon, Lyon", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Quartier des Marolles, Bruxelles  |  Et toute zone OSMnx valide", "level": 1, "size": 15, "gap": 0.52},
    {"text": "Protocole expérimental", "bold": True, "size": 16, "gap": 0.38},
    {"text": "N missions générées par ville (seedées, reproductibles)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "M politiques OR-Tools testées sur les mêmes instances", "level": 1, "size": 15, "gap": 0.35},
    {"text": "P passes de vérification pour le déterminisme", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Sortie : meilleure politique par ville + coût composite + delta vs baseline", "level": 1, "size": 15, "gap": 0.52},
    {"text": "make multi-city-benchmark", "size": 13, "color": C_LIGHT_BLUE, "gap": 0.38},
], icon="")

# ── SECTION 9 ─────────────────────────────────────────────────────────────────
add_section_divider(prs, "09", "Architecture & Démo",
                    "Stack technique et fonctionnalités de la plateforme")

add_two_col_slide(prs, "Architecture Technique",
    left_bullets=[
        {"text": "Backend — FastAPI (Python)", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  API REST + WebSocket", "size": 14, "color": C_GREY},
        {"text": "  OR-Tools (Google) : solveur VRPTW", "size": 14, "color": C_GREY},
        {"text": "  OSMnx + NetworkX : graphe routier", "size": 14, "color": C_GREY},
        {"text": "  Open-Meteo : météo temps réel", "size": 14, "color": C_GREY},
        {"text": "  SQLite : leaderboard persistant", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "Frontend — Next.js 14 + TypeScript", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Leaflet.js : carte interactive", "size": 14, "color": C_GREY},
        {"text": "  Visualisation des tournées", "size": 14, "color": C_GREY},
        {"text": "  Replay animé Humain vs IA", "size": 14, "color": C_GREY},
        {"text": "  Panneau de mission & sidebar", "size": 14, "color": C_GREY},
    ],
    right_bullets=[
        {"text": "Fonctionnalités clés", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Mode Joueur : l'humain choisit ses routes", "size": 14, "color": C_GREY},
        {"text": "  Mode IA : solveur OR-Tools joue seul", "size": 14, "color": C_GREY},
        {"text": "  Mode Versus : Humain vs IA en temps réel", "size": 14, "color": C_GREY},
        {"text": "  Auto-suggestion : IA conseille le joueur", "size": 14, "color": C_GREY},
        {"text": "  Incidents : reroutage automatique", "size": 14, "color": C_GREY},
        {"text": "  Météo réelle appliquée à la vitesse", "size": 14, "color": C_GREY},
        {"text": "", "size": 10},
        {"text": "DevOps & Qualité", "bold": True, "size": 16, "color": C_GOLD},
        {"text": "  Docker Compose (make docker)", "size": 14, "color": C_GREY},
        {"text": "  Suite de tests (pytest) : 15+ tests", "size": 14, "color": C_GREY},
        {"text": "  make repro-check : vérif. SHA-256", "size": 14, "color": C_GREY},
        {"text": "  Makefile : install, dev, test, docker", "size": 14, "color": C_GREY},
    ])

add_bullet_slide(prs, "Fonctionnalités Avancées", [
    {"text": "Debug solveur en temps réel", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "GET /api/missions/{id}/solver-debug → snapshot du contexte solveur", "level": 1, "size": 14, "color": C_LIGHT_BLUE, "gap": 0.35},
    {"text": "Décision de boost automatique selon nombre d'incidents, taille mission, météo", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Gestion des incidents en cours de mission", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "API de replanification POST /api/missions/{id}/replan-incident", "level": 1, "size": 14, "gap": 0.35},
    {"text": "Matrice de coût live recalculée avec les rues bloquées (core_data/live_time_matrix_incident.npy)", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Cache LRU sur les routes candidates", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Évite de recalculer les itinéraires déjà demandés → accélération mesurée", "level": 1, "size": 14, "gap": 0.48},
    {"text": "Système de leaderboard (Panthéon)", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Scores persistants SQLite, classement mondial par ville et profil", "level": 1, "size": 14, "gap": 0.35},
], icon="")

# ── SECTION 10 ────────────────────────────────────────────────────────────────
add_section_divider(prs, "10", "Conclusion & Perspectives",
                    "Bilan et ouvertures")

add_bullet_slide(prs, "Ce que nous avons accompli", [
    {"text": "Graphes & Open Data", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Graphe orienté pondéré réel (OSMnx) avec analyse complète de centralité et robustesse", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Données météo live (Open-Meteo), adresses géocodées, matrices de coût NumPy", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Algorithmique", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Dijkstra + A* + A* bidirectionnel pour le cheminement", "level": 1, "size": 15, "gap": 0.35},
    {"text": "VRPTW via OR-Tools avec GLS, SA, Tabu, ALNS, ILS", "level": 1, "size": 15, "gap": 0.35},
    {"text": "K-Means spatial custom (NumPy) pour la sectorisation", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Résultats concrets", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "-36.7% de temps de trajet, -36.1% de distance vs approche gloutonne humaine", "level": 1, "size": 15, "color": C_GREEN, "gap": 0.35},
    {"text": "Reproductibilité 100% (taux 1.0) sur toutes les politiques testées", "level": 1, "size": 15, "color": C_GREEN, "gap": 0.35},
], icon="")

add_bullet_slide(prs, "Perspectives & Ouvertures", [
    {"text": "Open Data enrichie", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Intégration des flux GTFS Île-de-France Mobilités → graphe multimodal (véhicule + métro)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Données de trafic temps réel (HERE, TomTom) pour pondérer les arêtes dynamiquement", "level": 1, "size": 15, "gap": 0.48},
    {"text": "Centralité avancée", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Eigenvector Centrality : identifier les nœuds bien connectés à d'autres nœuds importants", "level": 1, "size": 15, "gap": 0.35},
    {"text": "PageRank sur le graphe routier → prioriser les axes de transit", "level": 1, "size": 15, "gap": 0.48},
    {"text": "IA & Apprentissage", "bold": True, "size": 17, "color": C_GOLD, "gap": 0.45},
    {"text": "Renforcer le modèle avec un plan d'expérience complet (12-20 instances seedées)", "level": 1, "size": 15, "gap": 0.35},
    {"text": "Approche Reinforcement Learning pour remplacer l'heuristique gloutonne", "level": 1, "size": 15, "gap": 0.35},
], icon="")

# Final slide
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(slide, 0, 0, 13.33, 7.5, C_BG_DARK)
add_rect(slide, 0, 2.5, 13.33, 2.8, C_RED)
add_rect(slide, 0, 2.45, 13.33, 0.12, C_GOLD)
add_rect(slide, 0, 5.25, 13.33, 0.12, C_GOLD)
add_text(slide, "Merci", 0.5, 2.7, 12.3, 1.2,
         font_size=64, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Des questions ?",
         0.5, 3.9, 12.3, 0.7, font_size=22, color=C_GOLD, align=PP_ALIGN.CENTER)
add_text(slide, "Operation Noël  •  VRPTW + Open Data + Théorie des Graphes  •  2025-2026",
         0.5, 6.0, 12.3, 0.5, font_size=13, color=C_GREY, align=PP_ALIGN.CENTER)

output_path = "/home/bekkari/Documents/Graphes/Noel/rapport/Presentation_Operation_Noel.pptx"
import os
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
